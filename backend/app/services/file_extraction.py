"""
File extraction service for extracting text from various document formats
"""
import logging
from pathlib import Path
from typing import Optional
import io

logger = logging.getLogger(__name__)


def extract_file_content(file_content: bytes, filename: str) -> Optional[str]:
    """
    Extract text content from file bytes (for compatibility with existing code)
    
    Args:
        file_content: File content as bytes
        filename: Filename with extension
        
    Returns:
        Extracted text content or None if extraction fails
    """
    import tempfile
    import os
    
    # Determine file format from extension
    file_ext = filename.split('.')[-1].lower() if '.' in filename else ''
    
    # Create temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_ext}') as tmp_file:
        tmp_file.write(file_content)
        tmp_path = Path(tmp_file.name)
    
    try:
        # Use the async function (but call it synchronously for compatibility)
        import asyncio
        try:
            loop = asyncio.get_event_loop()
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        
        result = loop.run_until_complete(extract_text_from_file(tmp_path, file_ext))
        return result
    finally:
        # Clean up temporary file
        if tmp_path.exists():
            os.unlink(tmp_path)


async def extract_text_from_file(file_path: Path, file_format: str) -> Optional[str]:
    """
    Extract text content from a document file
    
    Args:
        file_path: Path to the file
        file_format: File format (pdf, docx, pptx, etc.)
        
    Returns:
        Extracted text content or None if extraction fails
    """
    if not file_path.exists():
        logger.error(f"File not found: {file_path}")
        return None
    
    file_format_lower = file_format.lower() if file_format else ""
    
    try:
        if file_format_lower == "pdf":
            return await _extract_from_pdf(file_path)
        elif file_format_lower in ["docx", "doc"]:
            return await _extract_from_docx(file_path)
        elif file_format_lower in ["pptx", "ppt"]:
            return await _extract_from_pptx(file_path)
        else:
            logger.warning(f"Unsupported file format for text extraction: {file_format}")
            return None
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {str(e)}", exc_info=True)
        return None


async def _extract_from_pdf(file_path: Path) -> Optional[str]:
    """Extract text from PDF file"""
    try:
        import PyPDF2
        
        text_parts = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                except Exception as e:
                    logger.warning(f"Error extracting text from PDF page {page_num}: {str(e)}")
                    continue
        
        return "\n\n".join(text_parts) if text_parts else None
    except ImportError:
        logger.error("PyPDF2 is not installed. Install it with: pip install PyPDF2")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return None


async def _extract_from_docx(file_path: Path) -> Optional[str]:
    """Extract text from DOCX file"""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n\n".join(text_parts) if text_parts else None
    except ImportError:
        logger.error("python-docx is not installed. Install it with: pip install python-docx")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        return None


async def _extract_from_pptx(file_path: Path) -> Optional[str]:
    """Extract text from PPTX file"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_parts.append(f"Slide {slide_num + 1}:\n" + "\n".join(slide_text))
        
        return "\n\n".join(text_parts) if text_parts else None
    except ImportError:
        logger.error("python-pptx is not installed. Install it with: pip install python-pptx")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        return None
